from pptx import Presentation

# создаем новую презентацию
prs = Presentation()
title_slide_layout = prs.slide_layouts[2]

text = [('choices', """choices(self, population, weights=None, *, cum_weights=None, k=1)
Return a k sized list of population elements chosen with replacement.  
If the relative weights or cumulative weights are not specified,
the selections are made with equal probability."""), ('expovariate', """expovariate(self, lambd)
Exponential distribution.
lambd is 1.0 divided by the desired mean.  It should be
nonzero.  (The parameter would be called "lambda", but that is
a reserved word in Python.)  Returned values range from 0 to
positive infinity if lambd is positive, and from negative
infinity to 0 if lambd is negative."""), ('gammavariate', """gammavariate(self, alpha, beta)
Gamma distribution.  Not the gamma function!
Conditions on the parameters are alpha > 0 and beta > 0.
The probability distribution function is:

           x ** (alpha - 1) * math.exp(-x / beta)
pdf(x) =  --------------------------------------
             math.gamma(alpha) * beta ** alpha"""), ('gauss', """gauss(self, mu, sigma)
Gaussian distribution.
mu is the mean, and sigma is the standard deviation.  This is
slightly faster than the normalvariate() function.
Not thread-safe without a lock around calls."""), ('lognormvariate', """lognormvariate(self, mu, sigma)
Log normal distribution.
If you take the natural logarithm of this distribution, you'll get a
normal distribution with mean mu and standard deviation sigma.
mu can have any value, and sigma must be greater than zero.""")]

for i in range(5):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = text[i][0]
    subtitle.text = text[i][1]

prs.save('res.pptx')
